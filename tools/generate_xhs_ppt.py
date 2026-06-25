"""Xiaohongshu showcase PPT — premium aesthetic edition"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

# ── Premium pastel palette ──
CREAM   = RGBColor(0xFB, 0xF7, 0xF2)
BLUSH   = RGBColor(0xF2, 0xCE, 0xD7)
ROSE    = RGBColor(0xD9, 0x7A, 0x8C)
MAUVE   = RGBColor(0xC4, 0x9B, 0xA8)
SAGE    = RGBColor(0xA7, 0xC0, 0xA4)
MINT    = RGBColor(0xBC, 0xD4, 0xC4)
WARM    = RGBColor(0xE8, 0xD5, 0xC0)
GOLD    = RGBColor(0xD4, 0xA5, 0x74)
SOFT_BG = RGBColor(0xFD, 0xF6, 0xF1)
BROWN   = RGBColor(0x5C, 0x4B, 0x47)
CHARCOAL= RGBColor(0x3D, 0x3A, 0x38)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)

# ── Helpers ──
def bg(slide, color=SOFT_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def txt(slide, l, t, w, h, text, size=18, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei', italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font; p.alignment = align
    if italic: p.font.italic = True
    return tf

def rect(slide, l, t, w, h, fill=None, border=None, radius=0.15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = radius / max(w, h) * 10 if max(w,h) > 0 else 0.1
    return shape

def circle(slide, l, t, d, fill):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.fill.background()
    return c

def add_paragraph(tf, text, size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT, spacing=6):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = 'Microsoft YaHei'; p.alignment = align
    p.space_before = Pt(spacing)
    return p

img_dir = os.path.join(os.path.dirname(__file__), '..', 'gallery', 'xiaohongshu')
imgs = {
    'ghibli': os.path.join(img_dir, '01_ghibli_hillside.jpg'),
    'cat':    os.path.join(img_dir, '02_cat_window.jpg'),
    'desk':   os.path.join(img_dir, '03_desk_setup.jpg'),
    'sakura': os.path.join(img_dir, '04_cherry_blossom.jpg'),
}

# ════════════════════════════
# SLIDE 1 — COVER (magazine editorial)
# ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, SOFT_BG)

# Decorative blob
c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-2), Inches(7), Inches(7))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xFD, 0xEA, 0xEF); c.line.fill.background()

c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(3.5), Inches(5), Inches(5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0xF0, 0xF6, 0xF1); c2.line.fill.background()

# Small accent dots
for (dx, dy) in [(8.5, 0.8), (1.2, 6.2), (9.1, 6.5)]:
    circle(s, dx, dy, 0.15, ROSE)

# Main title
txt(s, 1.0, 1.2, 8, 0.6, '✨', size=30, align=PP_ALIGN.CENTER)
txt(s, 1.0, 1.9, 8, 1.0, '我的 AI 画室', size=52, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)

# Decorative line
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.8), Inches(3.0), Inches(2.4), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ROSE; line.line.fill.background()

txt(s, 1.5, 3.3, 7, 0.7, '免费 · 12 秒出图 · 没有次数限制', size=20, color=MAUVE, align=PP_ALIGN.CENTER)
txt(s, 1.5, 4.3, 7, 0.6, '打字描述你想要的画面\n它帮你画出来', size=16, color=RGBColor(0x9E, 0x94, 0x90), align=PP_ALIGN.CENTER)

# Small corner text
txt(s, 1.5, 6.5, 7, 0.4, '往下翻看画作 👇', size=14, color=MAUVE, align=PP_ALIGN.CENTER)

# ════════════════════════════
# SLIDE 2 — THE ART (collage layout)
# ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, SOFT_BG)

# Title area
txt(s, 0.5, 0.2, 9, 0.5, '🖼️  我的 AI 画作', size=34, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
txt(s, 0.5, 0.7, 9, 0.4, '以下四张 · 全部免费生成 · 每张约 12 秒', size=14, color=MAUVE, align=PP_ALIGN.CENTER)

# 2×2 grid — larger squares, gallery-wall aesthetic
layout = [
    ('sakura', 0, 0, '🌸 樱花大道'),
    ('ghibli', 1, 0, '🌿 治愈山坡'),
    ('cat',    0, 1, '🐱 窗台小猫'),
    ('desk',   1, 1, '🩷 粉色桌面'),
]

sq = 2.5       # square image size
label_h = 0.45 # label below image
col_gap = 1.4  # horizontal gap between columns
row_gap = 0.6  # vertical gap between rows
total_w = sq * 2 + col_gap
total_h = (sq + label_h) * 2 + row_gap
start_x = (10 - total_w) / 2
start_y = 1.0

for key, col, row, label in layout:
    path = imgs[key]
    if not os.path.exists(path): continue

    fx = start_x + col * (sq + col_gap)
    fy = start_y + row * (sq + label_h + row_gap)

    # Polaroid card
    rect(s, fx - 0.1, fy - 0.1, sq + 0.2, sq + label_h + 0.2, fill=WHITE)
    # Image — square, no stretch
    s.shapes.add_picture(path, Inches(fx), Inches(fy), Inches(sq), Inches(sq))
    # Label
    txt(s, fx, fy + sq + 0.06, sq, label_h - 0.06, label,
        size=13, color=BROWN, align=PP_ALIGN.CENTER)

# ════════════════════════════
# SLIDE 3 — CAPABILITIES (cards with depth)
# ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, SOFT_BG)

txt(s, 0.5, 0.2, 9, 0.5, '🤖 一个工具全搞定', size=34, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
txt(s, 0.5, 0.7, 9, 0.4, '不用装一堆 App · 不用来回切换', size=14, color=MAUVE, align=PP_ALIGN.CENTER)

cards = [
    ('🎨', 'AI 绘画', '打字描述\n自动扩写\n12 秒出图', ROSE),
    ('👁️', '图片识别', '上传即分析\n论文/表格/照片\n全都能读', SAGE),
    ('💬', '智能对话', '回答问题\n写文章·翻译\n写代码·查资料', GOLD),
    ('📄', '网页阅读', '输入网址\n自动提取文字\n读论文神器', MAUVE),
]

for i, (emoji, title, desc, accent) in enumerate(cards):
    x = 0.5 + i * 2.35
    # Shadow card (offset)
    rect(s, x + 0.08, 1.6 + 0.08, 2.0, 4.8, fill=RGBColor(0xE8, 0xE0, 0xDC))
    # Main card
    rect(s, x, 1.6, 2.0, 4.8, fill=WHITE)

    # Accent circle at top
    circle(s, x + 0.55, 1.85, 0.8, accent)
    tf = s.shapes.add_textbox(Inches(x + 0.55), Inches(1.85), Inches(0.8), Inches(0.8))
    p = tf.text_frame.paragraphs[0]; p.text = emoji; p.font.size = Pt(22); p.alignment = PP_ALIGN.CENTER

    txt(s, x + 0.15, 2.85, 1.7, 0.45, title, size=20, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
    txt(s, x + 0.25, 3.55, 1.5, 2.5, desc, size=13, color=BROWN, align=PP_ALIGN.CENTER)

    # Bottom accent dot
    circle(s, x + 0.8, 6.05, 0.1, accent)

# ════════════════════════════
# SLIDE 4 — COST (elegant comparison)
# ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, SOFT_BG)

txt(s, 0.5, 0.2, 9, 0.5, '💰 费用 · 不玩文字游戏', size=34, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)

# Three columns
cols = [
    ('🆓', '永久免费\n3 项', ['AI 绘画', '图片识别', '网页抓取'], SAGE),
    ('💬', '按量付费\n1 项', ['智能对话 · 没月租\n重度用月约 20-30'], GOLD),
    ('📊', '对比一下', ['主流 AI 订阅\n每月 100-200 元\n这个是它的 1/5'], ROSE),
]

for i, (emoji, title, items, accent) in enumerate(cols):
    x = 0.8 + i * 3.1
    rect(s, x, 1.4, 2.6, 5.0, fill=WHITE)

    # Header
    hdr = rect(s, x, 1.4, 2.6, 1.1, fill=accent)
    tf = s.shapes.add_textbox(Inches(x), Inches(1.45), Inches(2.6), Inches(1.0))
    p = tf.text_frame.paragraphs[0]; p.text = f'{emoji}  {title}'; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER

    # Items
    for j, item in enumerate(items):
        txt(s, x + 0.2, 2.8 + j * 0.9, 2.2, 0.7, item, size=16, color=CHARCOAL, align=PP_ALIGN.CENTER)

# Bottom highlight
rect(s, 1.5, 6.7, 7, 0.5, fill=RGBColor(0xFD, 0xEA, 0xEF))
txt(s, 1.5, 6.72, 7, 0.45, '💡 画图免费 · 识图免费 · 对话便宜 · 用多少花多少', size=18, bold=True, color=ROSE, align=PP_ALIGN.CENTER)

# ════════════════════════════
# SLIDE 5 — CTA (clean invitation)
# ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, SOFT_BG)

# Large decorative circle
c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(1.0), Inches(6), Inches(6))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xFD, 0xEA, 0xEF); c.line.fill.background()

txt(s, 1.5, 2.0, 7, 0.8, '🙋', size=40, align=PP_ALIGN.CENTER)
txt(s, 1.5, 2.7, 7, 0.8, '想试试吗？', size=44, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)

# Divider
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.8), Inches(3.6), Inches(2.4), Inches(0.03))
line.fill.solid(); line.fill.fore_color.rgb = ROSE; line.line.fill.background()

txt(s, 1.5, 4.0, 7, 0.5, '自己动手 →  网上搜得到  免费开源', size=18, color=BROWN, align=PP_ALIGN.CENTER)
txt(s, 1.5, 4.6, 7, 0.5, '找人帮忙 →  装好配好  包教包会', size=18, color=BROWN, align=PP_ALIGN.CENTER)

txt(s, 1.5, 5.5, 7, 0.6, '感兴趣的评论区打个招呼', size=24, bold=True, color=ROSE, align=PP_ALIGN.CENTER)
txt(s, 1.5, 6.1, 7, 0.4, '我都会回复 📩', size=14, color=MAUVE, align=PP_ALIGN.CENTER)

# ── SAVE ──
out = os.path.join(os.path.dirname(__file__), '..', 'docs', 'xhs-showcase.pptx')
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
