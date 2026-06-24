"""Generate a professional showcase PPT for Claude Lab"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ──
BG_DARK  = RGBColor(0x0D, 0x11, 0x17)
BG_CARD  = RGBColor(0x16, 0x1B, 0x22)
ACCENT   = RGBColor(0x3B, 0x82, 0xF6)  # blue
ACCENT2  = RGBColor(0x10, 0xB9, 0x81)  # green
ACCENT3  = RGBColor(0xF5, 0x9E, 0x0B)  # amber
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x9C, 0xA3, 0xAF)
DARK_TXT = RGBColor(0x1F, 0x29, 0x37)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    return shape

def add_title_bar(slide):
    """Top accent line"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

# ═══════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)

add_text_box(slide, 1.5, 1.2, 10, 1.2, '🧪 Claude Lab', font_size=56, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.4, 10, 0.8, '一个终端，四个AI大脑', font_size=32, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.4, 10, 0.6, 'DeepSeek + Gemini + GPT-4o + Flux — 零锁定，全免费', font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# Divider line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.3), Inches(3.333), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, 1.5, 4.7, 10, 0.5, '开源 · 模型无关 · 一键安装 · 科研级', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.8, 10, 0.4, 'GitHub: hqw7484-maker/claude-lab', font_size=14, color=RGBColor(0x6B, 0x72, 0x80), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 2: PROBLEM
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)

add_text_box(slide, 1, 0.5, 11, 0.6, '🤔 现在的 AI 工具有什么问题？', font_size=36, bold=True, color=WHITE)

problems = [
    ('🔒', '厂商锁定', 'ChatGPT / Claude / Gemini 各自封闭\n换模型 = 数据迁移 + 重新适应 + 重复付费'),
    ('💰', '订阅疲劳', 'ChatGPT Plus $20/月 + Midjourney $10/月\n+ Gemini $20/月 = $50+/月'),
    ('🖼️', '功能割裂', '对话用 A，画图用 B，读图用 C\n窗口切来切去，上下文全丢'),
    ('🔧', '配置复杂', 'API Key、环境变量、命令行\n普通用户根本不知道怎么下手'),
]

for i, (emoji, title, desc) in enumerate(problems):
    x = 0.8 + i * 3.1
    add_card(slide, x, 1.5, 2.8, 4.5)
    add_text_box(slide, x + 0.3, 1.7, 2.2, 0.5, f'{emoji}  {title}', font_size=22, bold=True, color=WHITE)
    add_text_box(slide, x + 0.3, 2.4, 2.2, 3.2, desc, font_size=14, color=GRAY)

add_text_box(slide, 1, 6.4, 11, 0.5, '有没有一个工具，能一次性解决所有这些问题？', font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 3: SOLUTION
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)

add_text_box(slide, 1, 0.5, 11, 0.6, '✅ Claude Lab — 一站式解决方案', font_size=36, bold=True, color=WHITE)

solutions = [
    ('🔓 模型无关', '不绑定任何厂商\nDeepSeek / Gemini / GPT-4o 随时切换\n改一行配置，换一个脑子'),
    ('💸 极限省钱', '只有 DeepSeek 花钱（~1元/月）\n视觉识别 → Gemini 免费\nAI 绘画 → Flux 免费'),
    ('🧩 功能聚合', '一个终端搞定一切\n对话 · 画图 · 读图 · 查资料\n上下文全程保留，不中断'),
    ('⚡ 一键安装', '.\一键安装.ps1 回车即可\n自动检测环境、安装依赖\n引导填写 API Key'),
]

for i, (title, desc) in enumerate(solutions):
    x = 0.8 + i * 3.1
    add_card(slide, x, 1.5, 2.8, 4.8)
    add_text_box(slide, x + 0.3, 1.7, 2.2, 0.5, title, font_size=22, bold=True, color=ACCENT2)
    add_text_box(slide, x + 0.3, 2.4, 2.2, 3.5, desc, font_size=14, color=GRAY)

# ═══════════════════════════════════════
# SLIDE 4: AI ART DEMO
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)

add_text_box(slide, 1, 0.3, 11, 0.6, '🎨 AI 绘画 — 一句话出大片（免费）', font_size=36, bold=True, color=WHITE)

# Left: process
add_card(slide, 0.8, 1.2, 5.5, 5.5)
add_text_box(slide, 1.1, 1.4, 5, 0.4, '💬 用户输入', font_size=18, bold=True, color=ACCENT)
add_text_box(slide, 1.3, 1.9, 4.8, 0.5, 'gen "赛博朋克城市夜景，霓虹雨"', font_size=16, color=WHITE, font_name='Consolas')

add_text_box(slide, 1.1, 2.6, 5, 0.4, '🎬 AI 扩写（自动）', font_size=18, bold=True, color=ACCENT)
add_text_box(slide, 1.3, 3.1, 4.8, 1.2, 'PBR materials, cinematic volumetric lighting,\nOctane Render, 8K, photorealistic,\nwet asphalt reflections, neon reflections...', font_size=14, color=GRAY, font_name='Consolas')

add_text_box(slide, 1.1, 4.5, 5, 0.4, '⚙️ 轮询流水线', font_size=18, bold=True, color=ACCENT)
add_text_box(slide, 1.3, 5.0, 4.8, 1.5, 'Flux.1 (High-Res) → Turbo → Standard\n哪个先出用哪个\n✅ Flux.1 成功 · 12.05 秒\n💡 最终模型：Flux.1 (High-Res)', font_size=14, color=GRAY, font_name='Consolas')

# Right: image placeholder
add_card(slide, 7, 1.2, 5.5, 5.5)
add_text_box(slide, 7.5, 1.5, 4.5, 0.4, '🖼️ 生成结果', font_size=18, bold=True, color=ACCENT)

# Check if demo image exists and insert it
img_path = os.path.join(os.path.dirname(__file__), '..', 'gallery', 'art_2.jpg')
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(7.5), Inches(2.2), Inches(4.5), Inches(4.5))
else:
    add_text_box(slide, 8, 3.5, 4, 1, '[ Demo image placeholder ]', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 5: VISION DEMO
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)

add_text_box(slide, 1, 0.3, 11, 0.6, '👁️ AI 视觉 — 看图说话（免费）', font_size=36, bold=True, color=WHITE)

# Left: command
add_card(slide, 0.8, 1.2, 5.5, 2.3)
add_text_box(slide, 1.1, 1.4, 5, 0.4, '📥 命令', font_size=18, bold=True, color=ACCENT)
add_text_box(slide, 1.3, 1.9, 4.8, 1.2, 'vg art_2.jpg\n→ 调用 GPT-4o 高精度视觉解析', font_size=16, color=WHITE, font_name='Consolas')

# Right: free tags
add_card(slide, 7, 1.2, 5.5, 2.3)
add_text_box(slide, 7.3, 1.4, 5, 0.4, '🏷️ 可用模型', font_size=18, bold=True, color=ACCENT)
add_text_box(slide, 7.5, 1.9, 4.5, 1.2, 'v  <image>  →  Gemini 2.0 Flash（免费）\nvg <image>  →  GPT-4o（免费·GitHub Models）', font_size=16, color=WHITE, font_name='Consolas')

# Bottom: analysis result
add_card(slide, 0.8, 3.8, 11.7, 3.2)
add_text_box(slide, 1.1, 4.0, 11, 0.4, '📋 解析结果', font_size=18, bold=True, color=ACCENT)
add_text_box(slide, 1.3, 4.5, 10.8, 2.2,
    '赛博朋克风格城市夜景。路面因雨水湿滑，霓虹灯光在沥青上形成鲜艳倒影。两侧现代化建筑上\n'
    '布满明亮广告屏，部分可见日文和亚洲风格图案。流线型未来感汽车停在路边。远处天际线是\n'
    '密集高楼和无尽灯火。高饱和度霓虹灯配合积水反射，形成丰富电影级光影对比。',
    font_size=16, color=GRAY)

# ═══════════════════════════════════════
# SLIDE 6: BRAIN SWITCHING
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)

add_text_box(slide, 1, 0.3, 11, 0.6, '🧠 大脑切换 — 会话中热插拔，上下文不丢', font_size=36, bold=True, color=WHITE)

brains = [
    ('🧠', 'DeepSeek V4 Pro', 'ds', '逻辑与计算核心\n数学推理 / 代码生成\n主力日常使用', ACCENT),
    ('👁️', 'Gemini 2.0 Flash', 'gm', '视觉解析\n100万 token 长文本\n多模态分析', ACCENT2),
    ('🌐', 'GPT-4o', 'gp', '联网调研\n综合分析\n高精度识别', ACCENT3),
]

for i, (emoji, name, cmd, desc, accent) in enumerate(brains):
    x = 1.5 + i * 3.8
    add_card(slide, x, 1.5, 3.2, 4.5)
    # Accent top line on card
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(3.2), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.fill.background()

    add_text_box(slide, x + 0.3, 1.9, 2.6, 0.6, f'{emoji}  {name}', font_size=24, bold=True, color=WHITE)
    add_text_box(slide, x + 0.3, 2.7, 2.6, 0.4, f'命令: {cmd}', font_size=18, color=accent, font_name='Consolas')
    add_text_box(slide, x + 0.3, 3.4, 2.6, 2.2, desc, font_size=14, color=GRAY)

add_text_box(slide, 1, 6.5, 11, 0.4, '一句话切换，无需重启对话，无需导出数据 — 同一个会话，多个大脑协同工作', font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 7: ARCHITECTURE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)

add_text_box(slide, 1, 0.3, 11, 0.6, '🏗️ 架构 — 极简，但强大', font_size=36, bold=True, color=WHITE)

# ── Top: Core engine ──
core_w, core_h = 5.0, 1.4
core_x = (13.333 - core_w) / 2
core_y = 1.3

# Glow effect (larger semi-transparent rect behind)
glow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(core_x - 0.1), Inches(core_y - 0.1), Inches(core_w + 0.2), Inches(core_h + 0.2))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5C); glow.line.fill.background()

core = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(core_x), Inches(core_y), Inches(core_w), Inches(core_h))
core.fill.solid(); core.fill.fore_color.rgb = BG_CARD
core.line.color.rgb = ACCENT; core.line.width = Pt(1.5)

tf = core.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'Claude Code CLI'; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = 'DeepSeek V4 Pro  (Anthropic API 协议)'; p2.font.size = Pt(14); p2.font.color.rgb = ACCENT; p2.font.name = 'Microsoft YaHei'; p2.alignment = PP_ALIGN.CENTER

# ── Arrow: core → tools ──
arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(13.333/2 - 0.15), Inches(core_y + core_h + 0.1), Inches(0.3), Inches(0.5))
arrow1.fill.solid(); arrow1.fill.fore_color.rgb = ACCENT; arrow1.line.fill.background()

# ── Middle: 4 tool boxes ──
tools = [
    ('👁️  v.py', '视觉网关', 'Gemini + GPT-4o\n双引擎驱动', ACCENT),
    ('🎨  generate_art.py', 'AI 绘画', 'Flux → Turbo\n→ Standard 流水线', ACCENT2),
    ('🌐  fetch.py', '网页抓取', '零依赖\nPython 标准库', ACCENT3),
    ('📊  plot_sine.py', '数据可视化', 'Matplotlib\n图表生成', RGBColor(0xA8, 0x55, 0xF7)),
]

tool_y = 3.4
tool_w, tool_h = 2.7, 1.8
gap = 0.35
start_x = (13.333 - (tool_w * 4 + gap * 3)) / 2

for i, (name, title, desc, accent) in enumerate(tools):
    tx = start_x + i * (tool_w + gap)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(tx), Inches(tool_y), Inches(tool_w), Inches(tool_h))
    card.fill.solid(); card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = accent; card.line.width = Pt(1)

    # Accent top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(tx), Inches(tool_y), Inches(tool_w), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

    tf = card.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = accent; p.font.name = 'Consolas'; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(12); p2.font.color.rgb = GRAY; p2.font.name = 'Microsoft YaHei'; p2.alignment = PP_ALIGN.CENTER

# ── Arrow: tools → memory ──
arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(13.333/2 - 0.15), Inches(tool_y + tool_h + 0.1), Inches(0.3), Inches(0.4))
arrow2.fill.solid(); arrow2.fill.fore_color.rgb = ACCENT2; arrow2.line.fill.background()

# ── Bottom: Memory ──
mem_w, mem_h = 7.0, 1.0
mem_x = (13.333 - mem_w) / 2
mem_y = tool_y + tool_h + 0.7

mem = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(mem_x), Inches(mem_y), Inches(mem_w), Inches(mem_h))
mem.fill.solid(); mem.fill.fore_color.rgb = BG_CARD
mem.line.color.rgb = ACCENT2; mem.line.width = Pt(1)

tf = mem.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = '🧠  memory/  —  持久记忆系统'; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = 'rules.md (偏好)  ·  context.md (会话)  ·  MEMORY.md (索引)  —  关了重开，什么都记得'; p2.font.size = Pt(13); p2.font.color.rgb = GRAY; p2.font.name = 'Microsoft YaHei'; p2.alignment = PP_ALIGN.CENTER

# ── Annotation: cost badges ──
add_text_box(slide, 0.8, 6.8, 12, 0.4, '🆓 全部免费    💰 仅 DeepSeek 付费（≈1元/月）', font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 8: QUICK START
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)

add_text_box(slide, 1, 0.3, 11, 0.6, '⚡ 三步安装，全程不懂代码也能搞定', font_size=36, bold=True, color=WHITE)

steps = [
    ('1', '下载项目', 'git clone https://github.com/\nhqw7484-maker/claude-lab.git\ncd claude-lab', ACCENT),
    ('2', '配置密钥', '复制 .env.example → .env\n填入 DeepSeek API Key\n其他 Key 可选', ACCENT2),
    ('3', '一键安装启动', '.\\一键安装.ps1  → 自动装环境\n.\\start_lab.ps1 → 启动实验室', ACCENT3),
]

for i, (num, title, code, accent) in enumerate(steps):
    x = 1.2 + i * 3.9
    add_card(slide, x, 1.5, 3.4, 4.5)
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.2), Inches(1.7), Inches(0.9), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = 'Consolas'

    add_text_box(slide, x + 0.3, 2.9, 2.8, 0.5, title, font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, 3.5, 2.8, 2.2, code, font_size=14, color=GRAY, font_name='Consolas')

add_text_box(slide, 1, 6.5, 11, 0.4, '唯一硬性需求：DeepSeek API Key（platform.deepseek.com 注册即得）', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 9: COST
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)

add_text_box(slide, 1, 0.3, 11, 0.6, '💸 费用对比 — 比一杯奶茶还便宜', font_size=36, bold=True, color=WHITE)

# Comparison table
add_card(slide, 1.5, 1.5, 10, 4.5)

headers = ['功能', 'ChatGPT Plus', 'Claude Lab', '节省']
cols_x = [2.0, 4.5, 7.5, 10.0]
for h, cx in zip(headers, cols_x):
    add_text_box(slide, cx, 1.7, 2, 0.4, h, font_size=18, bold=True, color=ACCENT)

rows = [
    ('AI 对话',       '$20/月',  '≈ 1 元/月',   '99%'),
    ('AI 绘画',       '$10/月 (MJ)',  '免费 (Flux)',   '100%'),
    ('图片识别',      '需 Plus', '免费 (Gemini)',    '100%'),
    ('高精度识别',    '需 Plus', '免费 (GPT-4o)',    '100%'),
    ('网页抓取',      '需插件',  '免费 (内置)',    '100%'),
]

for j, (func, gpt, claude, save) in enumerate(rows):
    y = 2.4 + j * 0.65
    add_text_box(slide, 2.0, y, 2, 0.4, func, font_size=16, color=WHITE)
    add_text_box(slide, 4.5, y, 2.5, 0.4, gpt, font_size=16, color=RGBColor(0xEF, 0x44, 0x44))
    add_text_box(slide, 7.5, y, 2.5, 0.4, claude, font_size=16, color=ACCENT2)
    add_text_box(slide, 10.0, y, 1.5, 0.4, save, font_size=16, bold=True, color=ACCENT2)

add_text_box(slide, 1.5, 6.3, 10, 0.5, '每月总开销：$50+ (≈350元)  →  不到 5 元。省下 98%。', font_size=22, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 10: CTA
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)

add_text_box(slide, 1.5, 1.0, 10, 0.8, '🚀 开始你的 AI 实验室之旅', font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Divider
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(2.0), Inches(3.333), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, 1.5, 2.4, 10, 0.5, 'GitHub 开源 · MIT 协议 · 随便用随便改', font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.0, 10, 0.6, 'github.com/hqw7484-maker/claude-lab', font_size=28, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER, font_name='Consolas')

# Two CTAs
add_card(slide, 2.5, 4.2, 3.5, 2.2)
add_text_box(slide, 2.8, 4.4, 3, 0.4, '🆓 自己装', font_size=24, bold=True, color=WHITE)
add_text_box(slide, 2.8, 5.0, 3, 1.2, '.\一键安装.ps1 回车即可\n完全免费，MIT 开源', font_size=14, color=GRAY)

add_card(slide, 7.5, 4.2, 3.5, 2.2)
add_text_box(slide, 7.8, 4.4, 3, 0.4, '🔧 找我装', font_size=24, bold=True, color=WHITE)
add_text_box(slide, 7.8, 5.0, 3, 1.2, '68 元 · 30分钟 · 包教包会\nQQ: 3664049984', font_size=14, color=GRAY)

add_text_box(slide, 1.5, 6.8, 10, 0.4, '觉得有用？点个 Star ⭐ 让更多人看到', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ── SAVE ──
output_path = os.path.join(os.path.dirname(__file__), '..', 'docs', 'claude-lab-showcase.pptx')
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Slides: {len(prs.slides)}')
